import os
import fitz
import docx
import pptx
# import pickle
import time
from langchain_community.llms import Ollama
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.embeddings import SpacyEmbeddings
from langchain_community.vectorstores import FAISS
from langchain.chains.question_answering import load_qa_chain

def read_pdf(path):
    doc = fitz.open(path)
    text = ""
    for page in doc:
        text += page.get_text() + "\n"
    return text
  
def read_docx(path):
  doc = docx.Document(path)
  text = ""
  for para in doc.paragraphs:
      text += para.text + "\n"
  return text

def read_pptx(path):
  pres = pptx.Presentation(path)
  text = ""
  for slide in pres.slides:
    for shape in slide.shapes:
      if hasattr(shape, "text"):
        text += shape.text + "\n"
  return text

def read_txt(path):
  with open(path, "r") as file:
    text = file.read()
  return text

def extract_text(folder_path):
  texts = []
  for filename in os.listdir(folder_path):
    filepath = os.path.join(folder_path, filename)
    if filename.endswith('.pdf'):
      texts.append(read_pdf(filepath))
    elif filename.endswith('.docx'):
      texts.append(read_docx(filepath))
    elif filename.endswith('.txt'):
      texts.append(read_txt(filepath))
    elif filename.endswith('.pptx'):
      texts.append(read_pptx(filepath))
    
    texts.append("\n\n")

  return " ".join(texts)

def main():
  folder_path = "C:\\Users\\acer\\Documents\\docs" # Change this to your folder path
  combined_text = extract_text(folder_path)
  
  splitter = RecursiveCharacterTextSplitter(
              chunk_size=1000,
              chunk_overlap=100
              )
  chunks = splitter.split_text(text=combined_text)
  
  embeddings = SpacyEmbeddings()
  
  VectorStore = FAISS.from_texts(chunks, embedding=embeddings)
  
  # if os.path.exists(f"data.pkl"):
  #   with open(f"data.pkl", "rb") as f:
  #     VectorStore = pickle.load(f)
  # else:
  #   VectorStore = FAISS.from_texts(chunks, embedding=embeddings)
  #   with open(f"data.pkl", "wb") as f:
  #     pickle.dump(VectorStore, f)
  
  llm = Ollama(model="mistral")
  
  chain = load_qa_chain(llm=llm, chain_type="stuff")
  
  while True:
    query = input("Question: ")
    
    query = f"{query} - Anda harus menjawab dalam bahasa Indonesia"
    
    start_time = time.time()
    docs = VectorStore.similarity_search(query=query, k=5)
      
    response = chain.invoke({
          'input_documents': docs,
          'question': query
      })
    
    elapsed_time = time.time() - start_time
    minutes = int(elapsed_time // 60)
    seconds = elapsed_time % 60
    
    print("Answer:")
    print(response['output_text'])
    print(f"Waktu yang dibutuhkan: {minutes} menit {seconds:.2f} detik\n")
    
if __name__ == "__main__":
    main()
